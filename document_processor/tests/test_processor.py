"""Tests for document processor with mocked DB and OCR."""

from __future__ import annotations

import json
import uuid
from datetime import date
from unittest.mock import AsyncMock, MagicMock

import pytest

from document_processor.processor import (
    _build_chunk_metadata,
    process_document,
    strip_html,
)


def _make_row(
    doc_type: str = "agenda_item",
    content_raw: str | None = None,
    file_path: str | None = None,
    title: str = "Test Document",
    meeting_date: date | None = date(2024, 1, 15),
    committee_name: str = "Board of Directors",
    meeting_id: str = "mtg_001",
    agenda_item_id: str | None = "ai_001",
    source_url: str = "https://example.com/doc",
    metadata: dict | None = None,
) -> MagicMock:
    """Create a mock asyncpg.Record row."""
    row = MagicMock()
    row_data = {
        "id": uuid.uuid4(),
        "document_type": doc_type,
        "title": title,
        "content_raw": content_raw,
        "file_path": file_path,
        "meeting_date": meeting_date,
        "committee_name": committee_name,
        "meeting_id": meeting_id,
        "agenda_item_id": agenda_item_id,
        "source_url": source_url,
        "metadata": metadata or {},
    }
    row.__getitem__ = lambda self, key: row_data[key]
    row.get = lambda key, default=None: row_data.get(key, default)
    return row


class _AsyncContextManager:
    """Helper that wraps a value as an async context manager."""

    def __init__(self, value):
        self.value = value

    async def __aenter__(self):
        return self.value

    async def __aexit__(self, *args):
        return False


@pytest.fixture
def mock_pool():
    """Create a mock asyncpg connection pool.

    asyncpg's pool.acquire() and conn.transaction() both return objects
    usable directly as `async with ...` — they are NOT coroutines.
    We use MagicMock + _AsyncContextManager to replicate this.
    """
    conn = MagicMock()
    conn.execute = AsyncMock()
    conn.transaction.return_value = _AsyncContextManager(None)

    pool = MagicMock()
    pool.acquire.return_value = _AsyncContextManager(conn)
    pool._conn = conn
    return pool


@pytest.fixture
def mock_ocr():
    return AsyncMock()


class TestStripHTML:
    def test_strips_tags(self):
        html = "<p>Hello <strong>world</strong></p>"
        result = strip_html(html)
        assert "Hello" in result
        assert "world" in result
        assert "<p>" not in result
        assert "<strong>" not in result

    def test_preserves_text(self):
        assert strip_html("plain text") == "plain text"

    def test_handles_nested_divs(self, sample_html):
        text = strip_html(sample_html)
        assert "Board Resolution 2024-15" in text
        assert "$450 million" in text
        assert "<div>" not in text


class TestBuildChunkMetadata:
    def test_includes_all_fields(self):
        row = _make_row()
        meta = json.loads(_build_chunk_metadata(row))
        assert meta["meeting_date"] == "2024-01-15"
        assert meta["committee_name"] == "Board of Directors"
        assert meta["meeting_id"] == "mtg_001"
        assert meta["agenda_item_id"] == "ai_001"
        assert meta["title"] == "Test Document"
        assert meta["source_url"] == "https://example.com/doc"

    def test_skips_none_fields(self):
        row = _make_row(
            meeting_date=None,
            committee_name=None,
            agenda_item_id=None,
        )
        meta = json.loads(_build_chunk_metadata(row))
        assert "meeting_date" not in meta
        assert "committee_name" not in meta
        assert "agenda_item_id" not in meta
        assert "meeting_id" in meta


class TestProcessDocument:
    @pytest.mark.asyncio
    async def test_agenda_item_chunks_correctly(self, mock_pool, mock_ocr):
        """agenda_item with content_raw is HTML-stripped and chunked."""
        paragraph = "The board discussed budget allocations for the upcoming year. " * 30
        html = f"<div><p>{paragraph}</p></div>"
        row = _make_row(doc_type="agenda_item", content_raw=html)

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is True
        assert chunk_count >= 1
        mock_ocr.extract.assert_not_called()
        mock_pool._conn.execute.assert_called()

    @pytest.mark.asyncio
    async def test_missing_file_path_fails(self, mock_pool, mock_ocr):
        """Attachment without file_path fails cleanly."""
        row = _make_row(doc_type="attachment", file_path=None)

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is False
        assert chunk_count == 0
        mock_pool._conn.execute.assert_called()

    @pytest.mark.asyncio
    async def test_dry_run_no_db_writes(self, mock_pool, mock_ocr):
        """Dry run processes but doesn't write to DB."""
        paragraph = "Board meeting notes with important content. " * 30
        html = f"<p>{paragraph}</p>"
        row = _make_row(doc_type="agenda_item", content_raw=html)

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd", dry_run=True)

        assert success is True
        assert chunk_count >= 1
        mock_pool._conn.execute.assert_not_called()
        mock_pool.acquire.assert_not_called()

    @pytest.mark.asyncio
    async def test_failed_transaction_no_partial_chunks(self, mock_pool, mock_ocr):
        """If transaction fails, no partial chunks should remain."""
        paragraph = "Content for chunking in test document. " * 30
        html = f"<p>{paragraph}</p>"
        row = _make_row(doc_type="agenda_item", content_raw=html)

        mock_pool._conn.execute.side_effect = Exception("DB connection lost")

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is False
        assert chunk_count == 0

    @pytest.mark.asyncio
    async def test_metadata_propagated_to_chunks(self, mock_pool, mock_ocr):
        """Document metadata (meeting_date, committee, etc.) appears in chunk metadata."""
        paragraph = "Meeting content for metadata propagation test. " * 30
        html = f"<p>{paragraph}</p>"
        row = _make_row(
            doc_type="agenda_item",
            content_raw=html,
            meeting_date=date(2024, 3, 20),
            committee_name="Finance Committee",
            meeting_id="mtg_finance_001",
        )

        await process_document(row, mock_pool, mock_ocr, "kent_sd")

        calls = mock_pool._conn.execute.call_args_list
        chunk_inserts = [c for c in calls if "chunks" in str(c)]
        assert len(chunk_inserts) >= 1

        # conn.execute(SQL, tenant, doc_id, idx, content, tokens, metadata)
        # Positional args: [0]=SQL, [1]=tenant, ..., [6]=metadata
        meta_json = chunk_inserts[0][0][6]
        meta = json.loads(meta_json)
        assert meta["meeting_date"] == "2024-03-20"
        assert meta["committee_name"] == "Finance Committee"
        assert meta["meeting_id"] == "mtg_finance_001"

    @pytest.mark.asyncio
    async def test_ocr_called_for_attachment(self, mock_pool, mock_ocr, tmp_path):
        """Attachment documents should call OCR service."""
        fake_file = tmp_path / "test.pdf"
        fake_file.write_text("fake pdf")

        from document_processor.ocr_client import OCRResult

        mock_ocr.extract.return_value = OCRResult(
            text="Extracted text from PDF. " * 30,
            source="digital",
            page_count=2,
            char_count=750,
        )

        row = _make_row(doc_type="attachment", file_path=str(fake_file))

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is True
        mock_ocr.extract.assert_called_once_with(str(fake_file))

    @pytest.mark.asyncio
    async def test_pptx_file_routes_to_extractor(self, mock_pool, mock_ocr, tmp_path):
        """PowerPoint .pptx files bypass OCR and use pptx_extractor."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Board enrollment data for 2024. " * 30
        path = tmp_path / "enrollment.pptx"
        prs.save(str(path))

        row = _make_row(doc_type="attachment", file_path=str(path))

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is True
        assert chunk_count >= 1
        mock_ocr.extract.assert_not_called()

    @pytest.mark.asyncio
    async def test_ppsx_file_routes_to_extractor(self, mock_pool, mock_ocr, tmp_path):
        """PowerPoint .ppsx files also bypass OCR."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Budget presentation slides content. " * 30
        path = tmp_path / "budget.ppsx"
        prs.save(str(path))

        row = _make_row(doc_type="attachment", file_path=str(path))

        success, chunk_count = await process_document(row, mock_pool, mock_ocr, "kent_sd")

        assert success is True
        mock_ocr.extract.assert_not_called()

    @pytest.mark.asyncio
    async def test_pptx_extraction_tier_in_metadata(self, mock_pool, mock_ocr, tmp_path):
        """PowerPoint chunks should have extraction_tier: 0 in metadata."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "Policy review presentation content here. " * 30
        path = tmp_path / "policy.pptx"
        prs.save(str(path))

        row = _make_row(doc_type="attachment", file_path=str(path))

        await process_document(row, mock_pool, mock_ocr, "kent_sd")

        calls = mock_pool._conn.execute.call_args_list
        chunk_inserts = [c for c in calls if "chunks" in str(c)]
        assert len(chunk_inserts) >= 1

        meta_json = chunk_inserts[0][0][6]
        meta = json.loads(meta_json)
        assert meta["extraction_tier"] == 0

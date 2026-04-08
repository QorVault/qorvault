#!/usr/bin/env python3
"""One-off backfill: count embedded images in already-processed PDFs.

Reads every completed document with a file_path, counts images using
pypdf's page.images, and writes image_count into the metadata JSONB
column.  Does not re-extract text or re-chunk — metadata-only update.

Usage:
    cd /home/qorvault/projects/ksd-boarddocs-rag
    source .env   # or: set -a && source .env && set +a
    python3 scripts/backfill_image_counts.py [--dry-run] [--batch-size 200]
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import sys
import tempfile
import time
import zipfile
from pathlib import Path

import psycopg2
from dotenv import load_dotenv

# Suppress pypdf's noisy warnings about malformed PDFs
logging.getLogger("pypdf").setLevel(logging.ERROR)

load_dotenv(Path(__file__).resolve().parent.parent / ".env")

# Only open files under this directory — reject paths from DB that point elsewhere
ALLOWED_DATA_ROOT = "/home/qorvault/projects/ksd_forensic/"

# Max uncompressed size per zip entry when patching .ppsx (50 MB)
MAX_ZIP_ENTRY_BYTES = 50 * 1024 * 1024

FETCH_SQL = """
SELECT id, file_path, metadata
FROM documents
WHERE tenant_id = %s
  AND processing_status = 'complete'
  AND file_path IS NOT NULL
"""

UPDATE_SQL = """
UPDATE documents
SET metadata = metadata || %s::jsonb,
    updated_at = NOW()
WHERE id = %s
"""


def get_dsn() -> str:
    url = os.environ.get("DATABASE_URL")
    if url:
        return url
    host = os.environ.get("POSTGRES_HOST", "127.0.0.1")
    port = os.environ.get("POSTGRES_PORT", "5432")
    db = os.environ.get("POSTGRES_DB", "qorvault")
    user = os.environ.get("POSTGRES_USER", "qorvault")
    password = os.environ.get("POSTGRES_PASSWORD")
    if not password:
        print("ERROR: POSTGRES_PASSWORD not set. Source .env first.", file=sys.stderr)
        sys.exit(1)
    return f"postgresql://{user}:{password}@{host}:{port}/{db}"


def count_pdf_images(file_path: str) -> int:
    """Count embedded images in a PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    total = 0
    for page in reader.pages:
        try:
            total += len(page.images)
        except Exception:
            pass
    return total


_PPSX_CT = "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml"
_PPTX_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def _patch_ppsx_to_pptx(src: str, dst: str) -> None:
    """Rewrite [Content_Types].xml so python-pptx accepts .ppsx as .pptx."""
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            if item.file_size > MAX_ZIP_ENTRY_BYTES:
                raise ValueError(
                    f"Zip entry {item.filename} exceeds {MAX_ZIP_ENTRY_BYTES} bytes "
                    f"({item.file_size}), refusing to extract"
                )
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(_PPSX_CT.encode(), _PPTX_CT.encode())
            zout.writestr(item, data)


def count_pptx_images(file_path: str) -> int:
    """Count images in a PowerPoint file using python-pptx.

    For .ppsx files, patches [Content_Types].xml in a temp copy so
    python-pptx >= 1.0 accepts it.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    ext = os.path.splitext(file_path)[1].lower()
    tmp_path = None

    if ext in (".ppsx", ".pps"):
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        _patch_ppsx_to_pptx(file_path, tmp.name)
        tmp_path = tmp.name
        open_path = tmp_path
    else:
        open_path = file_path

    try:
        prs = Presentation(open_path)
        total = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total += 1
        return total
    finally:
        if tmp_path:
            os.unlink(tmp_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Backfill image_count for completed documents")
    parser.add_argument("--dry-run", action="store_true", help="Count but don't write to DB")
    parser.add_argument("--batch-size", type=int, default=200, help="Commit every N updates")
    parser.add_argument("--tenant", default="kent_sd")
    args = parser.parse_args()

    conn = psycopg2.connect(get_dsn())
    conn.autocommit = False
    cur = conn.cursor()

    cur.execute(FETCH_SQL, (args.tenant,))
    rows = cur.fetchall()
    print(f"Found {len(rows)} completed documents with file paths")

    updated = 0
    skipped = 0
    errors = 0
    with_images = 0
    batch_count = 0
    start = time.time()

    for doc_id, file_path, metadata in rows:
        # Skip if already backfilled
        if metadata and isinstance(metadata, dict) and "image_count" in metadata:
            skipped += 1
            continue
        if metadata and isinstance(metadata, str):
            try:
                meta = json.loads(metadata)
                if "image_count" in meta:
                    skipped += 1
                    continue
            except (json.JSONDecodeError, TypeError):
                pass

        if not file_path or "\x00" in file_path or not os.path.exists(file_path):
            skipped += 1
            continue

        # Path traversal guard: resolve symlinks so ../../etc/passwd and
        # symlink-based escapes are caught. Only open files under the
        # expected data root.
        resolved = os.path.realpath(file_path)
        if not resolved.startswith(ALLOWED_DATA_ROOT):
            errors += 1
            print(f"  BLOCKED path outside {ALLOWED_DATA_ROOT}: {file_path}")
            continue

        # Use resolved path for file operations — prevents TOCTOU race
        # where a symlink is swapped between validation and open.
        ext = os.path.splitext(resolved)[1].lower()
        try:
            if ext == ".pdf":
                image_count = count_pdf_images(resolved)
            elif ext in (".pptx", ".ppsx"):
                image_count = count_pptx_images(resolved)
            else:
                # Non-PDF/PPTX files (converted via LibreOffice) — can't
                # count images from the original format easily; mark as 0
                image_count = 0
        except Exception as e:
            errors += 1
            print(f"  ERROR {os.path.basename(file_path)}: {e}")
            continue

        image_count = max(image_count, 0)

        if image_count > 0:
            with_images += 1

        if not args.dry_run:
            cur.execute(UPDATE_SQL, (json.dumps({"image_count": image_count}), doc_id))
            batch_count += 1
            if batch_count >= args.batch_size:
                conn.commit()
                batch_count = 0

        updated += 1

        if updated % 500 == 0:
            elapsed = time.time() - start
            rate = updated / elapsed if elapsed > 0 else 0
            print(f"  Progress: {updated} processed, {with_images} with images, {rate:.0f}/s")

    if not args.dry_run and batch_count > 0:
        conn.commit()

    elapsed = time.time() - start
    cur.close()
    conn.close()

    print(f"\nDone in {elapsed:.1f}s")
    print(f"  Updated:      {updated}")
    print(f"  With images:  {with_images}")
    print(f"  Skipped:      {skipped}")
    print(f"  Errors:       {errors}")
    if args.dry_run:
        print("  (dry-run — no DB writes)")


if __name__ == "__main__":
    main()

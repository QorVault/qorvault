"""Shared fixtures for OCR service tests."""

from pathlib import Path

import pytest
from fpdf import FPDF


@pytest.fixture
def digital_pdf(tmp_path: Path) -> Path:
    """Create a PDF with embedded digital text."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    # Write enough text to pass MIN_DIGITAL_CHARS (100)
    text = "This is a test PDF with digital text content. " * 10
    pdf.multi_cell(0, 10, text)
    out = tmp_path / "digital.pdf"
    pdf.output(str(out))
    return out


@pytest.fixture
def minimal_pdf(tmp_path: Path) -> Path:
    """Create a PDF with very little text (below threshold)."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "Hi")
    out = tmp_path / "minimal.pdf"
    pdf.output(str(out))
    return out


@pytest.fixture
def empty_pdf(tmp_path: Path) -> Path:
    """Create a PDF with no text at all."""
    pdf = FPDF()
    pdf.add_page()
    out = tmp_path / "empty.pdf"
    pdf.output(str(out))
    return out


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """Create a simple .docx file."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Test document paragraph one.")
    doc.add_paragraph("Test document paragraph two.")
    out = tmp_path / "test.docx"
    doc.save(str(out))
    return out


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    """Create a simple .xlsx file."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 100])
    ws.append(["Beta", 200])
    out = tmp_path / "test.xlsx"
    wb.save(str(out))
    return out


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a simple .pptx file."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide Title"
    slide.placeholders[1].text = "Test slide body content"
    out = tmp_path / "test.pptx"
    prs.save(str(out))
    return out

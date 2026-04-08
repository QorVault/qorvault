"""Extract text from Office file formats (.docx, .pptx, .xlsx, etc.)."""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

OFFICE_EXTENSIONS = frozenset(
    {
        ".doc",
        ".docx",
        ".ppt",
        ".pptx",
        ".ppsx",
        ".xls",
        ".xlsx",
    }
)


def is_office_file(file_path: str) -> bool:
    return Path(file_path).suffix.lower() in OFFICE_EXTENSIONS


def extract_office(file_path: str) -> str:
    """Extract text from an Office file. Returns extracted text or empty string."""
    ext = Path(file_path).suffix.lower()

    if ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".doc":
        return _extract_doc(file_path)
    elif ext in (".pptx", ".ppsx"):
        return _extract_pptx(file_path)
    elif ext == ".ppt":
        return _extract_ppt(file_path)
    elif ext == ".xlsx":
        return _extract_xlsx(file_path)
    elif ext == ".xls":
        return _extract_xls(file_path)
    return ""


def _extract_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_doc(file_path: str) -> str:
    try:
        from docx import Document

        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.warning("Cannot parse .doc file %s: %s", Path(file_path).name, e)
        return ""


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def _extract_ppt(file_path: str) -> str:
    try:
        return _extract_pptx(file_path)
    except Exception as e:
        logger.warning("Cannot parse .ppt file %s: %s", Path(file_path).name, e)
        return ""


def _extract_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append("\t".join(cells))
    wb.close()
    return "\n".join(rows)


def _extract_xls(file_path: str) -> str:
    try:
        return _extract_xlsx(file_path)
    except Exception as e:
        logger.warning("Cannot parse .xls file %s: %s", Path(file_path).name, e)
        return ""

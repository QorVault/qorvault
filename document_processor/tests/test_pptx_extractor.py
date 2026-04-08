"""Tests for PowerPoint text extraction."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from document_processor.pptx_extractor import extract_pptx


def _make_pptx(tmp_path, slides_data: list[dict]) -> str:
    """Create a .pptx file from slide definitions.

    Each dict in slides_data can have:
      - texts: list[str]  — text boxes to add
      - notes: str         — speaker notes
      - table: list[list[str]] — rows × cols table data
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    for sd in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Text boxes
        for i, text in enumerate(sd.get("texts", [])):
            left = Inches(1)
            top = Inches(1 + i)
            txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.5))
            txBox.text_frame.text = text

        # Table
        if "table" in sd:
            rows_data = sd["table"]
            n_rows = len(rows_data)
            n_cols = len(rows_data[0]) if rows_data else 0
            if n_rows and n_cols:
                table_shape = slide.shapes.add_table(
                    n_rows,
                    n_cols,
                    Inches(1),
                    Inches(4),
                    Inches(6),
                    Inches(1),
                )
                table = table_shape.table
                for r, row_data in enumerate(rows_data):
                    for c, cell_text in enumerate(row_data):
                        table.cell(r, c).text = cell_text

        # Speaker notes
        if "notes" in sd:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = sd["notes"]

    path = str(tmp_path / "test.pptx")
    prs.save(path)
    return path


class TestExtractPptx:
    def test_extracts_shape_text(self, tmp_path):
        path = _make_pptx(
            tmp_path,
            [
                {"texts": ["Hello World", "Second paragraph"]},
            ],
        )
        text, count = extract_pptx(path)
        assert count == 1
        assert "Hello World" in text
        assert "Second paragraph" in text
        assert "--- Slide 1 ---" in text

    def test_extracts_speaker_notes(self, tmp_path):
        path = _make_pptx(
            tmp_path,
            [
                {"texts": ["Slide content"], "notes": "Important context here"},
            ],
        )
        text, count = extract_pptx(path)
        assert count == 1
        assert "Speaker notes: Important context here" in text

    def test_extracts_table_content(self, tmp_path):
        path = _make_pptx(
            tmp_path,
            [
                {
                    "texts": ["Table slide"],
                    "table": [
                        ["Year", "Enrollment"],
                        ["2023", "26000"],
                        ["2024", "25500"],
                    ],
                },
            ],
        )
        text, count = extract_pptx(path)
        assert count == 1
        assert "Year" in text
        assert "Enrollment" in text
        assert "26000" in text
        assert "25500" in text

    def test_empty_presentation(self, tmp_path):
        prs = Presentation()
        path = str(tmp_path / "empty.pptx")
        prs.save(path)

        text, count = extract_pptx(path)
        assert count == 0
        assert text == ""

    def test_slide_count_matches(self, tmp_path):
        path = _make_pptx(
            tmp_path,
            [
                {"texts": ["Slide 1"]},
                {"texts": ["Slide 2"]},
                {"texts": ["Slide 3"]},
            ],
        )
        text, count = extract_pptx(path)
        assert count == 3
        assert "--- Slide 1 ---" in text
        assert "--- Slide 2 ---" in text
        assert "--- Slide 3 ---" in text

    def test_multiple_slides_separated(self, tmp_path):
        path = _make_pptx(
            tmp_path,
            [
                {"texts": ["First slide content"]},
                {"texts": ["Second slide content"]},
            ],
        )
        text, count = extract_pptx(path)
        assert count == 2
        # Slides should be separated
        idx1 = text.index("First slide content")
        idx2 = text.index("Second slide content")
        assert idx1 < idx2

    def test_ppsx_extension(self, tmp_path):
        """Should work with .ppsx files too (same format)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        txBox.text_frame.text = "PPSX content"
        path = str(tmp_path / "show.ppsx")
        prs.save(path)

        text, count = extract_pptx(path)
        assert count == 1
        assert "PPSX content" in text

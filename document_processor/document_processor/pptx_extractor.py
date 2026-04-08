"""PowerPoint (.pptx/.ppsx) text extraction.

Extraction tier 0: native digital, no OCR required.
Extracts text from shapes, tables, and speaker notes.

python-pptx >= 1.0 rejects .ppsx files due to strict content-type
validation (slideshow.main+xml vs presentation.main+xml).  Since the
two formats share the same OpenXML ZIP structure, we work around this
by patching [Content_Types].xml inside the ZIP to swap the content type.
"""

from __future__ import annotations

import logging
import os
import tempfile
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


_PPSX_CT = "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml"
_PPTX_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"


def _patch_ppsx_to_pptx(src: str, dst: str) -> None:
    """Copy a .ppsx ZIP to *dst*, rewriting [Content_Types].xml so
    python-pptx accepts it as a .pptx.

    The only difference between .ppsx and .pptx is the content-type
    string in [Content_Types].xml — the slides/shapes are identical.
    """
    with zipfile.ZipFile(src, "r") as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(_PPSX_CT.encode(), _PPTX_CT.encode())
            zout.writestr(item, data)


def _open_presentation(file_path: str) -> tuple[Presentation, str | None]:
    """Open a PowerPoint file, handling .ppsx content-type issues.

    Returns (presentation, temp_path).  Caller should delete temp_path
    if not None.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in (".ppsx", ".pps"):
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        _patch_ppsx_to_pptx(file_path, tmp.name)
        try:
            prs = Presentation(tmp.name)
        except Exception:
            os.unlink(tmp.name)
            raise
        return prs, tmp.name

    return Presentation(file_path), None


def extract_pptx(file_path: str) -> tuple[str, int]:
    """Extract all text from a PowerPoint file.

    Returns (full_text, slide_count).  Each slide is separated by a
    ``--- Slide N ---`` header so the chunker can see boundaries.
    """
    prs, tmp_path = _open_presentation(file_path)
    try:
        return _extract_slides(prs, file_path)
    finally:
        if tmp_path:
            os.unlink(tmp_path)


def _extract_slides(prs: Presentation, file_path: str) -> tuple[str, int]:
    """Extract text from all slides in a presentation."""
    slides_text: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        # Shape text (text frames) and tables
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        parts.append(" | ".join(row_texts))
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            for paragraph in notes_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(f"Speaker notes: {text}")

        if parts:
            slides_text.append(f"--- Slide {slide_num} ---\n" + "\n".join(parts))

    slide_count = len(prs.slides)
    full_text = "\n\n".join(slides_text)

    logger.info(
        "Extracted %d slides from %s (%d chars)",
        slide_count,
        file_path,
        len(full_text),
    )
    return full_text, slide_count
